#!/usr/bin/env python3
"""Build the FastUrduGuard viva slide deck (formal structure, metric placeholders).

Usage (from repo root):
  python scripts/build_presentation.py
  python scripts/build_presentation.py --out path/to/deck.pptx
"""
from __future__ import annotations

import argparse
import io
from dataclasses import dataclass
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]


@dataclass(frozen=True)
class _FlowBox:
    x: float
    y: float
    w: float
    h: float
    text: str
    face: str
    edge: str = "#1e293b"
    alpha: float = 0.55

    def xc(self) -> float:
        return self.x + self.w / 2

    def yc(self) -> float:
        return self.y + self.h / 2

    def nn(self) -> tuple[float, float]:
        return self.xc(), self.y + self.h

    def ss(self) -> tuple[float, float]:
        return self.xc(), self.y

    def ee(self) -> tuple[float, float]:
        return self.x + self.w, self.yc()

    def ww(self) -> tuple[float, float]:
        return self.x, self.yc()


def _draw_flow_box(ax, box: _FlowBox) -> None:
    ax.add_patch(
        FancyBboxPatch(
            (box.x, box.y),
            box.w,
            box.h,
            boxstyle="round,pad=0.04",
            linewidth=1.75,
            edgecolor=box.edge,
            facecolor=box.face,
            alpha=box.alpha,
            zorder=4,
        )
    )
    ax.text(
        box.xc(),
        box.yc(),
        box.text,
        ha="center",
        va="center",
        fontsize=9.2,
        color="#0f172a",
        fontweight="600",
        linespacing=1.14,
        zorder=5,
    )


def _arrow(ax, p0: tuple[float, float], p1: tuple[float, float], **kw) -> None:
    props = dict(
        arrowstyle="-|>",
        mutation_scale=15,
        linewidth=2.0,
        color="#1e3a5f",
        shrinkA=3,
        shrinkB=3,
        zorder=3,
    )
    props.update(kw)
    ax.add_patch(FancyArrowPatch(p0, p1, **props))


def _wire_poly(ax, pts: list[tuple[float, float]], *, arrow_tail: bool = True) -> None:
    """Polylines in screen space; arrowhead only on the last segment when arrow_tail."""
    if len(pts) < 2:
        return
    if arrow_tail and len(pts) >= 2:
        if len(pts) > 2:
            for a, b in zip(pts[:-2], pts[1:-1]):
                ax.plot(
                    [a[0], b[0]],
                    [a[1], b[1]],
                    color="#1e3a5f",
                    lw=2.0,
                    solid_capstyle="round",
                    solid_joinstyle="round",
                    zorder=2,
                )
        _arrow(ax, pts[-2], pts[-1])
    else:
        for a, b in zip(pts[:-1], pts[1:]):
            ax.plot(
                [a[0], b[0]],
                [a[1], b[1]],
                color="#1e3a5f",
                lw=2.0,
                solid_capstyle="round",
                solid_joinstyle="round",
                zorder=2,
            )


def _placeholder_png(title: str, subtitle: str = "Replace before submission") -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(10.0, 5.2), dpi=120)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    fig.patch.set_facecolor("#f8fafc")
    ax.add_patch(
        FancyBboxPatch(
            (0.04, 0.08), 0.92, 0.84, boxstyle="round,pad=0.02",
            linewidth=2.0, edgecolor="#475569", facecolor="#f1f5f9", linestyle=(0, (5, 4)),
        )
    )
    ax.text(0.5, 0.58, title, ha="center", va="center", fontsize=19, fontweight="bold", color="#0f172a")
    ax.text(0.5, 0.38, subtitle, ha="center", va="center", fontsize=13, color="#475569")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


def _workflow_png() -> io.BytesIO:
    """System workflow: data pipeline → two ranks → Git → ensemble (schematic)."""
    fig_w, fig_h = 12.9, 6.1
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=130)
    ax.set_xlim(0, 13.05)
    ax.set_ylim(0, 6.52)
    ax.axis("off")
    fig.patch.set_facecolor("#fafafa")

    ax.add_patch(Rectangle((0.2, 3.45), 12.65, 2.93, facecolor="#ecfdf5",
                           edgecolor="none", zorder=0, alpha=0.88))
    ax.add_patch(Rectangle((0.2, 0.45), 12.65, 2.82, facecolor="#eff6ff",
                           edgecolor="none", zorder=0, alpha=0.9))

    ax.text(0.35, 6.22, "Data & NLP preprocessing", fontsize=11, fontweight="bold", color="#14532d")
    ax.text(0.35, 3.10, "PDC coordination", fontsize=11, fontweight="bold", color="#1e40af")

    w, h, g, y_top = 2.38, 1.05, 0.26, 5.12
    x0 = 0.54
    raw = _FlowBox(x0 + 0 * (w + g), y_top, w, h,
                   "Raw corpora\n(three public\nsources)", "#a7f3d0")
    unify = _FlowBox(x0 + 1 * (w + g), y_top, w, h,
                     "Unify +\ndeduplicate +\ntrain·val·test", "#a7f3d0")
    relabel = _FlowBox(x0 + 2 * (w + g), y_top, w, h,
                       "Four-class\nrelabel", "#a7f3d0")
    preprocess = _FlowBox(x0 + 3 * (w + g), y_top, w, h,
                          "Roman Urdu +\nparallel\npreprocess", "#86efac")

    tops = [raw, unify, relabel, preprocess]

    y_bot, h2, w2 = 1.16, 1.06, 2.42
    r0 = _FlowBox(1.12, y_bot, w2, h2,
                  "Rank 0\n(models A–C)", "#bfdbfe")
    r1 = _FlowBox(3.94, y_bot, w2, h2,
                  "Rank 1\n(models D–F)", "#bfdbfe")
    git_hub = _FlowBox(7.02, y_bot, 2.72, h2,
                       "Git + metrics +\naggregation", "#93c5fd")
    ens = _FlowBox(10.38, y_bot, w2, h2,
                   "Weighted\nensemble + explain +\nbenchmark", "#7dd3fc")

    for b in [*tops, r0, r1, git_hub, ens]:
        _draw_flow_box(ax, b)

    for left, right in zip(tops, tops[1:]):
        _arrow(ax, left.ee(), right.ww())

    bus_y = 3.88
    sx, sy = preprocess.ss()

    _wire_poly(ax, [(sx, sy), (sx, bus_y), (r0.xc(), bus_y), r0.nn()])
    _wire_poly(ax, [(sx, sy), (sx, bus_y), (r1.xc(), bus_y), r1.nn()])

    _arrow(ax, r0.ee(), git_hub.ww())
    _arrow(ax, r1.ee(), git_hub.ww())
    _arrow(ax, git_hub.ee(), ens.ww())

    ax.text(
        6.52,
        0.17,
        "Directed edges indicate processing / coordination flow (schematic — not timed to scale).",
        ha="center",
        fontsize=9.15,
        color="#64748b",
    )

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="#fafafa", pad_inches=0.12)
    plt.close(fig)
    buf.seek(0)
    return buf


def _set_title(slide, text: str) -> None:
    slide.shapes.title.text = text


def _body(slide):
    return slide.placeholders[1]


def _add_bullets(slide, title: str, lines: list[str]) -> None:
    _set_title(slide, title)
    tf = _body(slide).text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    t = prs.slides.add_slide(prs.slide_layouts[0])
    t.shapes.title.text = "FastUrduGuard"
    st = t.placeholders[1]
    st.text = (
        "Integrated parallel & distributed computing with NLP for Urdu misinformation detection\n\n"
        "Team: Sameer Asif (22I-0493) · Immad Shah (22I-0395)\n"
        "Course: Phase-3 · Cluster 1 (PDC + NLP)\n"
        "Repository: github.com/D3aThNdDeMiSe/FASTUrduGuard"
    )
    for p in st.text_frame.paragraphs:
        p.font.size = Pt(18)

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Problem statement",
        [
            "Urdu is widely used yet under-served by English-centric moderation and detection tools.",
            "False or misleading items spread quickly through text, memes, and Roman-Urdu social posts.",
            "Manual fact-checking does not scale; automated pipelines must handle noisy text and low-resource settings.",
            "Practical systems need fine-grained labels, explainability, and deployment-aware training — not lab-only accuracy.",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Literature review — baseline (Islam et al., 2025)",
        [
            'Islam et al., "Unified LLMs for Misinformation Detection in Low-Resource Linguistic Settings," arXiv:2506.01587, 2025.',
            "Introduces Urdu-LLD and strong binary {real, fake} results with six transformer encoders (e.g., mBERT, XLM-R).",
            "Single-workstation, sequential training; limited discussion of Roman Urdu, multi-class labels, or explainability.",
            "We adopt the same encoder families as a scientific control, but extend the problem setting and systems story.",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Research objectives",
        [
            "Unify and clean multiple public Urdu corpora under one schema with careful deduplication.",
            "Move beyond binary labels toward a four-way schema (Real, Fake, PartiallyFalse, Satire) with rule-based expansion.",
            "Fine-tune six encoders and combine them with a validation-weighted ensemble for robustness.",
            "Add explainability (attention + sentence-level LIME) and an inference-oriented evaluation path.",
            "Exploit two consumer GPUs on separate networks using a Git-coordinated training and aggregation workflow.",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Proposed solution — overview",
        [
            "Production-shaped pipeline: data fusion → multilingual preprocessing → distributed training manifests → metrics via Git.",
            "Optional Federated-Averaging of small LoRA adapters where bandwidth is constrained.",
            "Post-training: aggregator builds leaderboards/plots; benchmark and explanation scripts support demos.",
            "Technical stack: PyTorch, Hugging Face Transformers/Trainer, PEFT-LoRA, scikit-learn metrics.",
        ],
    )

    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(_workflow_png(), Inches(0.40), Inches(0.95), width=Inches(12.45))
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.20), Inches(12.0), Inches(0.65))
    tx.text_frame.text = "End-to-end workflow"
    tx.text_frame.paragraphs[0].font.size = Pt(28)
    tx.text_frame.paragraphs[0].font.bold = True

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Methodology — NLP",
        [
            "Unicode normalization (NFC), noise removal (URLs, HTML, elongated characters).",
            "Roman-Urdu detection and lexicon-aided transliteration when Roman tokens dominate.",
            "Stratified train / validation / test split on multi-class targets; multilingual tokenizers seq. length 256 (configurable).",
            "Fine-tuning: AdamW, weight decay, warmup + linear decay; optional LoRA on larger encoders.",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Methodology — PDC / coordination",
        [
            "Rank-specific model queues from coordinator/manifest.json (load-balanced by footprint).",
            "Home / campus networks behind NAT: asynchronous Git pushes of lightweight metrics and artefacts.",
            "GPU-serial scheduling per laptop (--placement_override gpu) avoids CPU contention from concurrent tokenization.",
            "Aggregation job merges checkpoints/metrics locally for ensemble scoring and reproducible plots.",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Challenges during implementation",
        [
            "Mixed-precision training: aligning Trainer AMP (fp16/bf16) with model dtypes to avoid scaler/clip failures.",
            "Consumer GPU stack: limited or absent BitsAndBytes sm_120 support — fallback to BF16/FP16 + LoRA instead of NF4 everywhere.",
            "Two-node parallelism without rendezvous TCP: GitOps choreography instead of classical NCCL collectives.",
            "Class imbalance across four labels vs. headline-heavy corpora; weighted ensemble mitigates weak tail classes.",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Comparison with the baseline approach",
        [
            "Labels: binary (Islam et al.) → four-way schema with explicit partial-fake and satire handling.",
            "Data: single curated head-line set → unified multi-source corpus with hash-based cross-corpus deduplication.",
            "Training: one machine sequential → two-rank manifest + Git-mediated coordination and optional FedAvg-of-LoRA.",
            "Analysis: accuracy-focused reporting → ensemble, confusion analysis, explainability hooks, and inference benchmarks (to be finalized).",
        ],
    )

    for title, ph, top in [
        ("Results — model & ensemble performance (placeholder)", "Insert bar chart: accuracy / macro-F1 / per-class F1 from aggregate.py", Inches(1.12)),
        ("Results — confusion matrix & error analysis (placeholder)", "Insert heatmap: test-set confusion matrix + optional calibration plot", Inches(1.12)),
    ]:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(_placeholder_png(ph), Inches(0.9), top, width=Inches(11.5))
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.0), Inches(0.68))
        tx.text_frame.text = title
        tx.text_frame.paragraphs[0].font.size = Pt(26)
        tx.text_frame.paragraphs[0].font.bold = True

    s = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s,
        "Technical readiness & oral examination",
        [
            "Laptop with conda/venv mirroring repository requirements.txt; datasets available locally.",
            "Trained checkpoints and parquet shards kept on-disk (large weights typically gitignored; metrics may be synced).",
            "Prepared live paths: scripts/node_train.py, aggregate pipeline, benchmark and explain_demo entry points.",
            "Thank you — questions welcome (Q&A guideline: approximately 5–7 minutes).",
        ],
    )
    for paragraph in _body(s).text_frame.paragraphs:
        paragraph.font.size = Pt(19)

    return prs


def default_output_paths() -> list[Path]:
    paths = [REPO_ROOT / "docs" / "FastUrduGuard_Presentation.pptx"]
    sibling = REPO_ROOT.parent / "Phase 3" / "FastUrduGuard_Presentation.pptx"
    if sibling.parent.is_dir():
        paths.append(sibling.resolve())
    return paths


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path, action="append", help="Destination .pptx (repeatable)")
    args = ap.parse_args()

    outs = list(args.out) if args.out else default_output_paths()
    prs = build_deck()
    for path in outs:
        path = path.resolve()
        path.parent.mkdir(parents=True, exist_ok=True)
        try:
            prs.save(str(path))
            print("Wrote", path)
        except PermissionError:
            print(
                "Skipped (file locked or open elsewhere):",
                path,
                "— close the deck in PowerPoint and re-run, or use --out.",
            )


if __name__ == "__main__":
    main()
